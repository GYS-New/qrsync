"""
ATALIAN — OYAK RENAULT için İO-GYS operasyonel sunumu oluşturur.

Çıktı: docs/sunum/ATALIAN_OYAK_RENAULT_GYS_Sunum.pptx
Görsel placeholder'ları: PowerPoint'te açıp ilgili kutuya kendi screenshot'ını
"Resim Ekle" ile yerleştirebilirsin (kutu otomatik silinir).
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── RENK PALETİ ────────────────────────────────────────────────────────────
ATALIAN_YESIL     = RGBColor(0x16, 0xA3, 0x4A)   # ana yeşil
ATALIAN_YESIL_K   = RGBColor(0x10, 0x6E, 0x32)   # koyu yeşil
ATALIAN_ACIK      = RGBColor(0xDC, 0xFC, 0xE7)   # açık yeşil arka
OYAK_SARI         = RGBColor(0xFB, 0xBF, 0x24)   # OYAK sarı vurgu
KOYU_GRI          = RGBColor(0x1F, 0x29, 0x37)
ORTA_GRI          = RGBColor(0x6B, 0x72, 0x80)
ACIK_GRI          = RGBColor(0xF3, 0xF4, 0xF6)
BEYAZ             = RGBColor(0xFF, 0xFF, 0xFF)
MAVI              = RGBColor(0x1D, 0x4E, 0xD8)
PLACEHOLDER_BG    = RGBColor(0xF8, 0xFA, 0xFC)
PLACEHOLDER_KEN   = RGBColor(0x94, 0xA3, 0xB8)
KIRMIZI           = RGBColor(0xDC, 0x26, 0x26)

# ── SUNUM BOYUT (16:9 widescreen) ──────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── YARDIMCI FONKSİYONLAR ──────────────────────────────────────────────────

def add_slide(arka_renk=BEYAZ):
    s = prs.slides.add_slide(BLANK)
    # arka plan dolgu
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = arka_renk
    bg.shadow.inherit = False
    return s

def add_box(slide, x, y, w, h, fill=None, line=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        if line_w is not None:
            sh.line.width = line_w
    sh.shadow.inherit = False
    return sh

def add_text(slide, x, y, w, h, text, *, font_size=14, bold=False, color=KOYU_GRI,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top  = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.name = font_name
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def add_header_band(slide, title, subtitle=None):
    """Slayt üst bandı: yeşil çubuk + başlık + alt başlık + Atalian/OYAK rozeti."""
    add_box(slide, 0, 0, prs.slide_width, Inches(0.85), fill=ATALIAN_YESIL_K)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(9), Inches(0.5),
             title, font_size=24, bold=True, color=BEYAZ)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.55), Inches(9), Inches(0.3),
                 subtitle, font_size=11, color=ATALIAN_ACIK)
    # sağ üstte rozet
    add_text(slide, Inches(10.0), Inches(0.2), Inches(3.0), Inches(0.5),
             "ATALIAN — OYAK RENAULT", font_size=11, bold=True, color=OYAK_SARI,
             align=PP_ALIGN.RIGHT)

def add_footer(slide, slide_no, total):
    """Alt çubuk: sayfa numarası + sistem adı."""
    add_box(slide, 0, prs.slide_height - Inches(0.35), prs.slide_width, Inches(0.35),
            fill=ACIK_GRI)
    add_text(slide, Inches(0.5), prs.slide_height - Inches(0.32), Inches(8), Inches(0.3),
             "İO-GYS — Görev Yönetim Sistemi", font_size=9, color=ORTA_GRI)
    add_text(slide, Inches(10.5), prs.slide_height - Inches(0.32), Inches(2.5), Inches(0.3),
             f"{slide_no} / {total}", font_size=9, color=ORTA_GRI, align=PP_ALIGN.RIGHT)

def add_image_placeholder(slide, x, y, w, h, label):
    """Bir görsel için checkered placeholder kutusu."""
    box = add_box(slide, x, y, w, h, fill=PLACEHOLDER_BG, line=PLACEHOLDER_KEN,
                  line_w=Pt(1.5))
    # Çapraz çizgiler (diagonal X)
    diag_color = RGBColor(0xCB, 0xD5, 0xE1)
    l1 = slide.shapes.add_connector(1, x, y, x + w, y + h)
    l1.line.color.rgb = diag_color
    l1.line.width = Pt(0.75)
    l2 = slide.shapes.add_connector(1, x + w, y, x, y + h)
    l2.line.color.rgb = diag_color
    l2.line.width = Pt(0.75)
    # Üst label
    add_text(slide, x, y + h/2 - Inches(0.25), w, Inches(0.3),
             "🖼  GÖRSEL", font_size=12, bold=True, color=ORTA_GRI, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + h/2 + Inches(0.05), w, Inches(0.3),
             label, font_size=11, color=PLACEHOLDER_KEN, align=PP_ALIGN.CENTER)

def add_bullet_list(slide, x, y, w, h, items, *, font_size=14, color=KOYU_GRI, bullet="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"{bullet}  {it}"
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
        r.font.name = 'Calibri'

def add_kpi_card(slide, x, y, w, h, baslik, deger, alt=None, renk=ATALIAN_YESIL):
    add_box(slide, x, y, w, h, fill=BEYAZ, line=renk, line_w=Pt(1.5))
    add_text(slide, x, y + Inches(0.1), w, Inches(0.3), baslik,
             font_size=10, bold=True, color=ORTA_GRI, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.4), w, Inches(0.7), deger,
             font_size=28, bold=True, color=renk, align=PP_ALIGN.CENTER)
    if alt:
        add_text(slide, x, y + h - Inches(0.4), w, Inches(0.3), alt,
                 font_size=9, color=ORTA_GRI, align=PP_ALIGN.CENTER)

# ── SAYFALAR ───────────────────────────────────────────────────────────────

SLIDES = []  # her sayfayı önce planla, sonra footer için toplam saysıyı bil

def s_kapak():
    s = add_slide(arka_renk=KOYU_GRI)
    # Sol yan yeşil aksent
    add_box(s, 0, 0, Inches(0.4), prs.slide_height, fill=ATALIAN_YESIL)
    # Üst rozet
    add_text(s, Inches(1), Inches(0.7), Inches(11), Inches(0.4),
             "ATALIAN  ×  OYAK RENAULT", font_size=14, bold=True,
             color=OYAK_SARI, align=PP_ALIGN.LEFT)
    # Ana başlık
    add_text(s, Inches(1), Inches(2.2), Inches(11), Inches(1.0),
             "İO-GYS", font_size=72, bold=True, color=BEYAZ)
    add_text(s, Inches(1), Inches(3.3), Inches(11), Inches(0.6),
             "Görev Yönetim Sistemi", font_size=28, color=ATALIAN_ACIK)
    add_text(s, Inches(1), Inches(4.1), Inches(11), Inches(0.5),
             "Frekansiyel Temizlik & Kontrol Operasyonu", font_size=18,
             color=ORTA_GRI)
    # Alt bilgi
    add_text(s, Inches(1), Inches(6.3), Inches(11), Inches(0.4),
             "Operasyonel Sunum — Haziran 2026", font_size=14,
             color=BEYAZ)
    add_text(s, Inches(1), Inches(6.7), Inches(11), Inches(0.4),
             "Web Paneli  •  Mobil Uygulama  •  Raporlama",
             font_size=12, color=ATALIAN_ACIK)
    return s

def s_icindekiler():
    s = add_slide()
    add_header_band(s, "İçindekiler", "Sunumun yapısı")
    items = [
        "1.  Sistem Genel Bakış — neden, ne, nasıl",
        "2.  Web Paneli — yetki seviyeleri ve ana sayfalar",
        "3.  Gösterge Paneli & Canlı Akış",
        "4.  Görev & Kural Yönetimi",
        "5.  Mobil Uygulama — personel akışı",
        "6.  QR / NFC Okutma & Süre Takibi",
        "7.  Ekstra Görev (Frekans Dışı)",
        "8.  Raporlama: Genel Rapor",
        "9.  Personel Değerlendirme",
        "10. Hakediş & Faturalama",
        "11. Güvenlik & Denetim",
        "12. Müşteri için Faydalar",
        "13. Yol Haritası ve Soru/Cevap",
    ]
    add_bullet_list(s, Inches(1), Inches(1.4), Inches(11), Inches(5.4),
                    items, font_size=18, bullet="")
    return s

def s_sistem_genel():
    s = add_slide()
    add_header_band(s, "1. Sistem Genel Bakış", "Neden İO-GYS?")
    # Sol kutu — Problem
    add_box(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(2.6),
            fill=ACIK_GRI, line=KIRMIZI, line_w=Pt(1))
    add_text(s, Inches(0.7), Inches(1.4), Inches(5.6), Inches(0.4),
             "⚠  ÖNCEDEN", font_size=14, bold=True, color=KIRMIZI)
    add_bullet_list(s, Inches(0.7), Inches(1.85), Inches(5.6), Inches(2),
                    [
                        "Kağıt/Excel ile manuel takip",
                        "Personel kontrolü zayıf — kim ne yaptı belirsiz",
                        "Müşteri denetimde 'ispatlı veri' yok",
                        "Fatura kalemleri tartışmalı",
                        "Ekstra işler hiç kayıt altına alınmıyor",
                    ], font_size=13)
    # Sağ kutu — Çözüm
    add_box(s, Inches(6.85), Inches(1.3), Inches(6.0), Inches(2.6),
            fill=ATALIAN_ACIK, line=ATALIAN_YESIL, line_w=Pt(1))
    add_text(s, Inches(7.05), Inches(1.4), Inches(5.6), Inches(0.4),
             "✓  İO-GYS İLE", font_size=14, bold=True, color=ATALIAN_YESIL_K)
    add_bullet_list(s, Inches(7.05), Inches(1.85), Inches(5.6), Inches(2),
                    [
                        "Otomatik kural-tabanlı görev üretimi",
                        "QR/NFC ile fiziksel doğrulama (manipüle edilemez)",
                        "Saniye hassasiyetiyle süre takibi",
                        "Müşteri kendi panelinden anlık denetim",
                        "Aylık hakediş raporları otomatik üretilir",
                    ], font_size=13)
    # Alt akış diyagramı
    add_text(s, Inches(0.5), Inches(4.2), Inches(12), Inches(0.4),
             "Çalışma Akışı", font_size=16, bold=True, color=KOYU_GRI)
    bx_w, bx_h = Inches(2.2), Inches(1.3)
    bx_y = Inches(4.8)
    steps = [
        ("📋", "Kural Tanımı", "TA web'den\nfrekans + saat\nkuralı girer"),
        ("⚙", "Otomatik Üretim", "Sistem her gün\nvardiyalara göre\ngörev üretir"),
        ("📱", "Mobil Okutma", "Personel lokasyon\nQR'ını okutup\nbaşlatır/tamamlar"),
        ("📊", "Anlık Rapor", "Web + müşteri paneli\ngörev başarısını\ncanlı gösterir"),
        ("💰", "Hakediş", "Ay sonu fatura\nkalemleri otomatik\nüretilir"),
    ]
    sx = Inches(0.5)
    for i, (icon, baslik, aciklama) in enumerate(steps):
        x = sx + (bx_w + Inches(0.15)) * i
        add_box(s, x, bx_y, bx_w, bx_h, fill=BEYAZ, line=ATALIAN_YESIL, line_w=Pt(1.2))
        add_text(s, x, bx_y + Inches(0.05), bx_w, Inches(0.4), icon,
                 font_size=22, align=PP_ALIGN.CENTER)
        add_text(s, x, bx_y + Inches(0.45), bx_w, Inches(0.3), baslik,
                 font_size=12, bold=True, color=ATALIAN_YESIL_K, align=PP_ALIGN.CENTER)
        add_text(s, x, bx_y + Inches(0.75), bx_w, Inches(0.55), aciklama,
                 font_size=9, color=ORTA_GRI, align=PP_ALIGN.CENTER)
    return s

def s_web_yetki():
    s = add_slide()
    add_header_band(s, "2. Web Paneli — Yetki Seviyeleri",
                    "Kim ne görür, kim ne yapabilir?")
    # 3 sütun
    cols = [
        ("SÜPER ADMİN", "Sistem sağlayıcı",
         ["Tüm firmalar", "Sistem yönetimi", "Lisans & ayarlar",
          "Audit log erişimi", "Yedekleme & arşiv"], MAVI),
        ("TENANT ADMIN", "Hizmet sağlayıcı yöneticisi (Atalian)",
         ["Bir firmanın tüm projeleri", "Personel atama",
          "Görev & lokasyon yönetimi", "Tüm raporlar", "Birim fiyat tanımları"],
         ATALIAN_YESIL),
        ("MÜŞTERİ KULLANICISI", "Denetim (OYAK Renault)",
         ["Yetkili lokasyonlar (read-only)", "Görev takip + canlı akış",
          "Kendi rapor + Excel indirme", "Çeklist sonuçları",
          "Personel performansı"], OYAK_SARI),
    ]
    col_w = Inches(4.0)
    col_h = Inches(5.2)
    sx = Inches(0.6)
    for i, (baslik, alt, items, renk) in enumerate(cols):
        x = sx + (col_w + Inches(0.25)) * i
        add_box(s, x, Inches(1.4), col_w, col_h, fill=BEYAZ, line=renk, line_w=Pt(2))
        # Üst başlık bandı
        add_box(s, x, Inches(1.4), col_w, Inches(0.7), fill=renk)
        add_text(s, x, Inches(1.5), col_w, Inches(0.5), baslik,
                 font_size=15, bold=True, color=BEYAZ, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(2.2), col_w, Inches(0.3), alt,
                 font_size=10, color=ORTA_GRI, align=PP_ALIGN.CENTER)
        add_bullet_list(s, x + Inches(0.2), Inches(2.7), col_w - Inches(0.4),
                        col_h - Inches(1.4), items, font_size=11, bullet="✓")
    return s

def s_web_dashboard():
    s = add_slide()
    add_header_band(s, "3. Gösterge Paneli — Vardiya KPI'ları",
                    "Anlık operasyon görünümü")
    # Sol: 3 KPI kartı
    add_kpi_card(s, Inches(0.5), Inches(1.4), Inches(2.4), Inches(1.5),
                 "VARDİYA 1 (23:30–07:30)", "%98", "180 toplam · 152 tamam · 1 kayıp")
    add_kpi_card(s, Inches(3.0), Inches(1.4), Inches(2.4), Inches(1.5),
                 "VARDİYA 2 (07:30–15:30)", "%34", "358 toplam · 121 tamam",
                 renk=OYAK_SARI)
    add_kpi_card(s, Inches(5.5), Inches(1.4), Inches(2.4), Inches(1.5),
                 "VARDİYA 3 (15:30–23:30)", "%0", "281 toplam · bugün başlamadı",
                 renk=ORTA_GRI)
    # Sağ üst: özellik listesi
    add_text(s, Inches(8.2), Inches(1.4), Inches(5), Inches(0.4),
             "Sayfa içeriği", font_size=14, bold=True, color=ATALIAN_YESIL_K)
    add_bullet_list(s, Inches(8.2), Inches(1.75), Inches(5), Inches(2),
                    ["Vardiya başına başarı yüzdesi",
                     "Toplam / tamamlanan / sapma / kayıp",
                     "Anlık personel aktivitesi (bildirim çubuğu)",
                     "Canlı görev akışı linki"], font_size=12)
    # Alt: görsel placeholder
    add_image_placeholder(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(3.7),
                          "Web — Gösterge Paneli (vardiya kartları + KPI)")
    return s

def s_web_canli():
    s = add_slide()
    add_header_band(s, "4. Canlı Görev Akışı",
                    "Saniye saniye operasyon takibi")
    # Sol: özellikler
    add_text(s, Inches(0.5), Inches(1.4), Inches(6), Inches(0.4),
             "Ne sağlar?", font_size=16, bold=True, color=ATALIAN_YESIL_K)
    add_bullet_list(s, Inches(0.5), Inches(1.85), Inches(6), Inches(4),
                    ["Son N saatin tüm görev hareketi",
                     "8 KPI kartı (Toplam / Tamamlandı / İşlemde / Beklemede / İptal / Gecikmeli / Zamanı Geçmiş / Ekstra)",
                     "Tıklanabilir filtre (kart üzerinde toggle)",
                     "Üst lokasyon filtresi (departman)",
                     "Otomatik 5 saniyede bir yenilenir",
                     "Yeni gelen kayıt sarı parlamayla vurgulanır",
                     "Çekliste ve çeklist sonucuna anlık erişim",
                     "İptal sebebi popup (kim, neden iptal etti)"],
                    font_size=12)
    # Sağ: placeholder
    add_image_placeholder(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.5),
                          "Web — Canlı Akış sayfası (KPI kartları + liste)")
    return s

def s_web_kurallar():
    s = add_slide()
    add_header_band(s, "5. Görev Kuralı & Lokasyon Yönetimi",
                    "Frekansiyel görev üretimi nasıl planlanır?")
    # Akış diyagramı
    add_text(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
             "Bir kuralın hayat döngüsü", font_size=14, bold=True, color=KOYU_GRI)
    steps = [
        ("Lokasyon Tanımı", "Üst lokasyon\n+ alt lokasyon\n+ QR/NFC etiket"),
        ("Kural Oluşturma", "Tanım + frekans\n+ vardiya + min/max\nsüre + atanan rol"),
        ("Otomatik Üretim", "Her gece kural\nbazlı görev kaydı\nproduce edilir"),
        ("Personel Atama", "Vardiya saatinde\nlokasyona göre\notomatik atama"),
        ("Mobile Tarama", "Personel QR ile\nbaşlatır/tamamlar\nveya çeklist doldurur"),
    ]
    bx_w, bx_h = Inches(2.3), Inches(1.4)
    bx_y = Inches(1.9)
    sx = Inches(0.5)
    for i, (b, a) in enumerate(steps):
        x = sx + (bx_w + Inches(0.15)) * i
        add_box(s, x, bx_y, bx_w, bx_h, fill=ATALIAN_ACIK, line=ATALIAN_YESIL, line_w=Pt(1))
        add_text(s, x, bx_y + Inches(0.15), bx_w, Inches(0.4), f"{i+1}",
                 font_size=18, bold=True, color=ATALIAN_YESIL_K, align=PP_ALIGN.CENTER)
        add_text(s, x, bx_y + Inches(0.5), bx_w, Inches(0.35), b,
                 font_size=11, bold=True, color=KOYU_GRI, align=PP_ALIGN.CENTER)
        add_text(s, x, bx_y + Inches(0.85), bx_w, Inches(0.55), a,
                 font_size=9, color=ORTA_GRI, align=PP_ALIGN.CENTER)
    # Görsel: lokasyon ekranı placeholder
    add_image_placeholder(s, Inches(0.5), Inches(3.7), Inches(6.0), Inches(3.2),
                          "Web — Lokasyon listesi + QR kod yönetimi")
    add_image_placeholder(s, Inches(6.8), Inches(3.7), Inches(6.0), Inches(3.2),
                          "Web — Görev Kuralı tanım formu")
    return s

def s_mobil_akis():
    s = add_slide()
    add_header_band(s, "6. Mobil Uygulama — Personel Akışı",
                    "Sahadaki kullanıcı deneyimi")
    steps = [
        ("Giriş", "Firma kodu\n+ şifre", "🔐"),
        ("Mesai Başı", "İş başı QR\nokutması", "⏱"),
        ("QR/NFC Okut", "Lokasyon\netiketini okut", "📷"),
        ("Başlat", "Görevi başlat\n(opsiyonel)", "▶"),
        ("İşi Yap", "Sahada görev\ngerçekleştirilir", "🧹"),
        ("Tamamla", "Tekrar QR okut\n+ Tamamla", "✓"),
        ("Mesai Çıkışı", "İş çıkışı QR\nokutması", "🏁"),
    ]
    bx_w, bx_h = Inches(1.7), Inches(2.0)
    bx_y = Inches(1.5)
    sx = Inches(0.4)
    for i, (b, a, icon) in enumerate(steps):
        x = sx + (bx_w + Inches(0.07)) * i
        add_box(s, x, bx_y, bx_w, bx_h, fill=BEYAZ, line=ATALIAN_YESIL, line_w=Pt(1.5))
        add_text(s, x, bx_y + Inches(0.15), bx_w, Inches(0.5), icon,
                 font_size=28, align=PP_ALIGN.CENTER)
        add_text(s, x, bx_y + Inches(0.7), bx_w, Inches(0.35), f"{i+1}. {b}",
                 font_size=11, bold=True, color=ATALIAN_YESIL_K, align=PP_ALIGN.CENTER)
        add_text(s, x, bx_y + Inches(1.1), bx_w, Inches(0.8), a,
                 font_size=9, color=ORTA_GRI, align=PP_ALIGN.CENTER)
        # ok
        if i < len(steps) - 1:
            arrow_x = x + bx_w + Emu(50000)
            add_text(s, arrow_x, bx_y + bx_h/2 - Inches(0.15), Inches(0.1),
                     Inches(0.3), "▶", font_size=11, color=ATALIAN_YESIL,
                     align=PP_ALIGN.CENTER)
    # Alt: 3 mobil ekran placeholder
    add_image_placeholder(s, Inches(0.5), Inches(3.9), Inches(4.0), Inches(3.0),
                          "Mobil — Ana Sayfa")
    add_image_placeholder(s, Inches(4.7), Inches(3.9), Inches(4.0), Inches(3.0),
                          "Mobil — QR Okutma Sonrası")
    add_image_placeholder(s, Inches(8.9), Inches(3.9), Inches(4.0), Inches(3.0),
                          "Mobil — Aktif Görev / Tamamla")
    return s

def s_qr_sure():
    s = add_slide()
    add_header_band(s, "7. QR/NFC Okutma & Süre Takibi",
                    "Sahteciliğe kapalı, denetlenebilir veri")
    # Üstte 3 ikon kart
    cards = [
        ("📍", "Fiziksel Doğrulama",
         "Personel lokasyonun yanındaki QR'ı / NFC etiketini\nokutmadan görev yapamaz."),
        ("⏱", "Min Süre",
         "Lokasyona göre tanımlanan en kısa çalışma süresi\ndolmadan tamamlama engellenir."),
        ("🚨", "Max Süre",
         "Üst sınırı aşan görevler sistem tarafından\n'Görev Zaman Aşımı' ile iptal edilir."),
    ]
    cx, cy = Inches(0.5), Inches(1.4)
    cw, ch = Inches(4.1), Inches(2.0)
    for i, (icon, b, a) in enumerate(cards):
        x = cx + (cw + Inches(0.13)) * i
        add_box(s, x, cy, cw, ch, fill=ATALIAN_ACIK, line=ATALIAN_YESIL, line_w=Pt(1.2))
        add_text(s, x, cy + Inches(0.1), cw, Inches(0.5), icon,
                 font_size=26, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(0.6), cw, Inches(0.4), b,
                 font_size=14, bold=True, color=ATALIAN_YESIL_K, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(1.05), cw, Inches(0.9), a,
                 font_size=11, color=KOYU_GRI, align=PP_ALIGN.CENTER)
    # Alt: defense in depth
    add_text(s, Inches(0.5), Inches(3.7), Inches(12), Inches(0.4),
             "Çift Katmanlı Koruma (Defense in Depth)",
             font_size=15, bold=True, color=KOYU_GRI)
    add_box(s, Inches(0.5), Inches(4.15), Inches(6.0), Inches(2.4),
            fill=BEYAZ, line=ATALIAN_YESIL, line_w=Pt(1))
    add_text(s, Inches(0.7), Inches(4.25), Inches(5.6), Inches(0.4),
             "1. Mobil katman", font_size=13, bold=True, color=ATALIAN_YESIL_K)
    add_bullet_list(s, Inches(0.7), Inches(4.65), Inches(5.6), Inches(1.8),
                    ["UI'da Tamamla butonu min süre dolana kadar disabled",
                     "Wall-clock (sistem saati) ile süre kontrolü",
                     "Bypass denemesi log'lanır"], font_size=11)
    add_box(s, Inches(6.85), Inches(4.15), Inches(6.0), Inches(2.4),
            fill=BEYAZ, line=MAVI, line_w=Pt(1))
    add_text(s, Inches(7.05), Inches(4.25), Inches(5.6), Inches(0.4),
             "2. Backend katman", font_size=13, bold=True, color=MAVI)
    add_bullet_list(s, Inches(7.05), Inches(4.65), Inches(5.6), Inches(1.8),
                    ["Sunucu zamanı esas alınır (manipüle edilemez)",
                     "Min süre dolmadan tamamlama → 400 reddedilir",
                     "Audit log: her bypass denemesi kayda alınır"], font_size=11)
    return s

def s_ekstra_gorev():
    s = add_slide()
    add_header_band(s, "8. Ekstra Görev — Frekans Dışı Çalışma",
                    "Planlı kuralın dışında yapılan iş")
    add_text(s, Inches(0.5), Inches(1.3), Inches(8), Inches(0.4),
             "Senaryo", font_size=14, bold=True, color=KOYU_GRI)
    add_text(s, Inches(0.5), Inches(1.7), Inches(8), Inches(0.8),
             "Lokasyonda öngörülmeyen bir durum oluşur (örn: yağ döküldü, müdahale gerekti). "
             "Personel kural fazlası bu işi sisteme kaydeder, gerekçesi ile birlikte.",
             font_size=12, color=KOYU_GRI)
    # Akış
    add_text(s, Inches(0.5), Inches(2.7), Inches(12), Inches(0.4),
             "Mobil ekstra görev akışı (v2)", font_size=13, bold=True, color=ATALIAN_YESIL_K)
    items = [
        ("1️⃣  QR Okut", "Personel lokasyon QR'ını okutur, 'Ekstra Görev Yap' butonu görünür."),
        ("2️⃣  Tanım & Gerekçe", "Görev tanımı dropdown'dan seçilir + gerekçe yazılır (min 10 karakter)."),
        ("3️⃣  Başlat", "'Başlat' butonu — kronometre devreye girer, görev İŞLEMDE durumuna geçer."),
        ("4️⃣  İşi Yap", "Kullanıcı işi sahada yapar, ana sayfaya dönebilir."),
        ("5️⃣  Tamamla", "Görev tamamlandığında 'Tamamla' tuşu. Backend süreyi otomatik hesaplar."),
    ]
    for i, (b, a) in enumerate(items):
        y = Inches(3.1 + 0.55 * i)
        add_text(s, Inches(0.6), y, Inches(2.5), Inches(0.4), b,
                 font_size=12, bold=True, color=KOYU_GRI)
        add_text(s, Inches(3.2), y, Inches(9.5), Inches(0.4), a,
                 font_size=11, color=KOYU_GRI)
    # Alt: raporlama
    add_box(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.85),
            fill=OYAK_SARI)
    add_text(s, Inches(0.7), Inches(6.18), Inches(12), Inches(0.3),
             "Raporlamada", font_size=11, bold=True, color=KOYU_GRI)
    add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.4),
             "Genel Rapor → 'Frekans Dışı Çalışmalar' sekmesinde personel, tarih, "
             "süre, GEREKÇE ile listelenir. Atalian/OYAK ay sonu denetiminde her "
             "ekstra işin nedenini görür.",
             font_size=10, color=KOYU_GRI)
    return s

def s_rapor_genel():
    s = add_slide()
    add_header_band(s, "9. Genel Rapor — Aylık Performans",
                    "Müşterinin gördüğü ana denetim raporu")
    # KPI kartları
    add_kpi_card(s, Inches(0.5), Inches(1.4), Inches(2.0), Inches(1.4),
                 "HEDEF", "10,856", "Kural-tabanlı görev")
    add_kpi_card(s, Inches(2.7), Inches(1.4), Inches(2.0), Inches(1.4),
                 "TAMAMLANAN", "9,247", "Başarı: %85", renk=ATALIAN_YESIL)
    add_kpi_card(s, Inches(4.9), Inches(1.4), Inches(2.0), Inches(1.4),
                 "SAPMA", "892", "Geç tamamlanan", renk=OYAK_SARI)
    add_kpi_card(s, Inches(7.1), Inches(1.4), Inches(2.0), Inches(1.4),
                 "KAYIP", "717", "Hiç yapılmamış", renk=KIRMIZI)
    add_kpi_card(s, Inches(9.3), Inches(1.4), Inches(2.0), Inches(1.4),
                 "EKSTRA", "143", "Gerekçeli", renk=MAVI)
    # Alt: sekmeler + görsel
    add_text(s, Inches(0.5), Inches(3.1), Inches(12), Inches(0.4),
             "Rapor sekmeleri", font_size=13, bold=True, color=KOYU_GRI)
    add_bullet_list(s, Inches(0.5), Inches(3.5), Inches(6), Inches(3),
                    ["Özet & Grafikler — pasta, çubuk grafikler",
                     "Grup Analizi — lokasyon grubu bazlı dağılım",
                     "Tamamlananlar — tüm tamamlanan görev kayıtları",
                     "Sapmalar — geç yapılanlar (ZAMANINDA_YAPILAMAYAN)",
                     "Kayıp Frekanslar — hiç yapılmamış (vardiya bitti, iptal vs.)",
                     "Frekans Dışı — ekstra görevler (gerekçeli)",
                     "Atanan Frekanslar — kim hangi göreve atanmıştı?"],
                    font_size=11)
    add_image_placeholder(s, Inches(6.8), Inches(3.4), Inches(6.0), Inches(3.4),
                          "Web — Genel Rapor sekmeleri")
    return s

def s_rapor_personel():
    s = add_slide()
    add_header_band(s, "10. Personel Değerlendirme",
                    "Kim ne kadar çalıştı, başarı kategorisi")
    # Sol: özellikler
    add_text(s, Inches(0.5), Inches(1.4), Inches(6), Inches(0.4),
             "Sayfada bulunan veriler", font_size=14, bold=True, color=ATALIAN_YESIL_K)
    add_bullet_list(s, Inches(0.5), Inches(1.85), Inches(6), Inches(4),
                    ["Tamamlanan görev sayısı",
                     "İptal sayısı",
                     "Ortalama görev süresi (per task)",
                     "Aktif gün sayısı",
                     "Günlük görev süresi (tüm tamamlamaların toplamı / aktif gün)",
                     "Başarı kategorisi: BAŞARILI / NORMAL / YETERSİZ / BAŞARISIZ",
                     "Cihaz eşleşme durumu",
                     "Filtreler: tarih, vardiya, üst lokasyon, personel"],
                    font_size=12)
    # Sağ: kategori legend
    add_text(s, Inches(6.85), Inches(1.4), Inches(6), Inches(0.4),
             "Başarı kategorileri (günlük çalışma)", font_size=13, bold=True,
             color=KOYU_GRI)
    cats = [
        ("0–1 saat", "BAŞARISIZ", KIRMIZI),
        ("1–3 saat", "YETERSİZ", OYAK_SARI),
        ("3–6 saat", "NORMAL", MAVI),
        ("6+ saat",  "BAŞARILI", ATALIAN_YESIL),
    ]
    for i, (sure, kat, renk) in enumerate(cats):
        y = Inches(1.95 + 0.55 * i)
        add_box(s, Inches(6.85), y, Inches(1.6), Inches(0.45), fill=renk)
        add_text(s, Inches(6.85), y + Inches(0.07), Inches(1.6), Inches(0.3),
                 sure, font_size=11, bold=True, color=BEYAZ, align=PP_ALIGN.CENTER)
        add_text(s, Inches(8.55), y + Inches(0.07), Inches(4), Inches(0.3),
                 kat, font_size=14, bold=True, color=renk)
    # Alt: placeholder
    add_image_placeholder(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.4),
                          "Web — Personel Değerlendirme tablosu")
    return s

def s_hakedis():
    s = add_slide()
    add_header_band(s, "11. Hakediş & Fatura",
                    "Birim fiyat × tamamlanan = otomatik fatura kalemi")
    # Sol: nasıl çalışır
    add_text(s, Inches(0.5), Inches(1.4), Inches(6), Inches(0.4),
             "Hakediş nasıl hesaplanır?", font_size=14, bold=True,
             color=ATALIAN_YESIL_K)
    add_bullet_list(s, Inches(0.5), Inches(1.85), Inches(6), Inches(4),
                    ["Her lokasyon × görev tanımı için birim fiyat tanımlanır",
                     "Tamamlanan görev sayısı × birim fiyat = kalem tutarı",
                     "Sapma ve kayıp görevler fatura dışıdır",
                     "Ekstra görevler ayrı kalem olarak gösterilir",
                     "Aylık/dönemsel olarak Excel export edilir",
                     "Müşteri kendi panelinden indirebilir"],
                    font_size=12)
    # Sağ: örnek tablo
    add_text(s, Inches(6.85), Inches(1.4), Inches(6), Inches(0.4),
             "Örnek hakediş satırı", font_size=13, bold=True, color=KOYU_GRI)
    headers = ["LOKASYON", "GÖREV", "ADET", "B.FİYAT", "TOPLAM"]
    rows = [
        ["DISGS", "WC Temizliği", "245", "8,50 ₺", "2.082,50 ₺"],
        ["DISGS", "Çay Bölgesi",  "62",  "12,00 ₺","744,00 ₺"],
        ["MONTAJ","Plastik Toplama","412","6,25 ₺","2.575,00 ₺"],
        ["MONTAJ","Dolly Transferi","78","18,00 ₺","1.404,00 ₺"],
        ["",      "TOPLAM",       "",    "",       "6.805,50 ₺"],
    ]
    cw = [Inches(1.2), Inches(1.6), Inches(0.8), Inches(1.0), Inches(1.3)]
    tx = Inches(6.85)
    ty = Inches(1.85)
    # Header
    for i, h in enumerate(headers):
        x = tx + sum(cw[:i], Emu(0))
        add_box(s, x, ty, cw[i], Inches(0.35), fill=ATALIAN_YESIL_K)
        add_text(s, x, ty + Inches(0.04), cw[i], Inches(0.27), h,
                 font_size=9, bold=True, color=BEYAZ, align=PP_ALIGN.CENTER)
    # Rows
    for ri, row in enumerate(rows):
        y = ty + Inches(0.35 + 0.4 * ri)
        bg = OYAK_SARI if ri == len(rows)-1 else (ACIK_GRI if ri % 2 else BEYAZ)
        for ci, val in enumerate(row):
            x = tx + sum(cw[:ci], Emu(0))
            add_box(s, x, y, cw[ci], Inches(0.4), fill=bg, line=PLACEHOLDER_KEN, line_w=Pt(0.5))
            bold = ri == len(rows)-1
            add_text(s, x, y + Inches(0.06), cw[ci], Inches(0.28), val,
                     font_size=10, bold=bold, color=KOYU_GRI, align=PP_ALIGN.CENTER)
    return s

def s_guvenlik():
    s = add_slide()
    add_header_band(s, "12. Güvenlik & Denetim",
                    "Veri bütünlüğü ve manipülasyonu önleme")
    items = [
        ("🔒", "QR/NFC Fiziksel Doğrulama",
         "Personel lokasyonun yanında olmadan görev yapamaz. QR/NFC etiketi her lokasyonda kalıcı."),
        ("⏱", "Sunucu Saatli Süre Hesabı",
         "Wall-clock = backend now() - baslatilma_tarihi. Mobil saatinin manipüle edilmesi etkisiz."),
        ("📝", "Audit Log",
         "Login, görev iptal, ekstra kayıt, min süre bypass denemesi vb. tüm kritik işlemler izlenir."),
        ("🚪", "Yetki Sınırı",
         "Müşteri kullanıcısı sadece kendi yetkili lokasyonlarını görür. Cross-firma erişim engellenir."),
        ("🛡", "Defense in Depth",
         "Mobil (1. katman) + backend (2. katman) validasyon — biri devre dışı kalsa diğer korur."),
        ("🔐", "KVKK Uyum",
         "Personel verisi minimum kapsamda tutulur. Audit log uzun süreli saklanır, ham veri arşivlenir."),
    ]
    bx_w, bx_h = Inches(6.1), Inches(1.5)
    sx, sy = Inches(0.5), Inches(1.4)
    for i, (icon, b, a) in enumerate(items):
        col = i % 2
        row = i // 2
        x = sx + (bx_w + Inches(0.15)) * col
        y = sy + (bx_h + Inches(0.15)) * row
        add_box(s, x, y, bx_w, bx_h, fill=BEYAZ, line=ATALIAN_YESIL, line_w=Pt(1))
        add_text(s, x + Inches(0.15), y + Inches(0.1), Inches(0.5), Inches(0.5),
                 icon, font_size=22)
        add_text(s, x + Inches(0.8), y + Inches(0.15), Inches(5), Inches(0.4),
                 b, font_size=12, bold=True, color=ATALIAN_YESIL_K)
        add_text(s, x + Inches(0.8), y + Inches(0.55), bx_w - Inches(1), Inches(0.9),
                 a, font_size=10, color=KOYU_GRI)
    return s

def s_musteri_faydalari():
    s = add_slide()
    add_header_band(s, "13. OYAK RENAULT İçin Faydalar",
                    "Müşterinin elde ettiği değer")
    cards = [
        ("📊", "Şeffaflık",
         "Her görevin ne zaman, kim tarafından, nerede yapıldığı dakika hassasiyetiyle görünür."),
        ("✅", "İspatlı Denetim",
         "QR/NFC ile fiziksel doğrulama — 'yapıldı dendi ama yapılmadı' tartışması ortadan kalkar."),
        ("💰", "Doğru Faturalama",
         "Hakediş otomatik üretilir; tartışmalı kalemler engellenir, ay sonu kapanış hızlanır."),
        ("🎯", "Performans Görünürlüğü",
         "Personel ve lokasyon bazında başarı oranı; problemli alanlar erken tespit edilir."),
        ("📱", "Anlık Erişim",
         "Web panelinden 7/24 erişim, mobil özet, Excel/CSV export — toplantı öncesi rapor hazır."),
        ("📈", "Sürekli İyileştirme",
         "Sapma ve kayıp analizi → hangi vardiyada hangi lokasyon sorunlu, root-cause hızlı bulunur."),
    ]
    bx_w, bx_h = Inches(4.05), Inches(1.7)
    sx, sy = Inches(0.5), Inches(1.4)
    for i, (icon, b, a) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = sx + (bx_w + Inches(0.1)) * col
        y = sy + (bx_h + Inches(0.15)) * row
        add_box(s, x, y, bx_w, bx_h, fill=ATALIAN_ACIK, line=ATALIAN_YESIL_K, line_w=Pt(1.5))
        add_text(s, x, y + Inches(0.1), bx_w, Inches(0.5), icon,
                 font_size=26, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.55), bx_w, Inches(0.4), b,
                 font_size=14, bold=True, color=ATALIAN_YESIL_K, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), y + Inches(0.95), bx_w - Inches(0.3), Inches(0.7),
                 a, font_size=10, color=KOYU_GRI, align=PP_ALIGN.CENTER)
    return s

def s_yol_haritasi():
    s = add_slide()
    add_header_band(s, "14. Yol Haritası — Sonraki 3 Ay",
                    "Geliştirme planı (Haziran–Ağustos 2026)")
    items = [
        ("Haziran", "Ekstra Görev V2 yayını",
         "Mobil 1.0.28 — gerekçe + süre takibi. Min/Max süre backend validasyonu."),
        ("Haziran", "Personel Değerlendirme iyileştirme",
         "Günlük görev süresi şeffaflığı, sıralama, sütun filtreleri."),
        ("Temmuz", "Mobil Hata Log Paneli",
         "Web panelde mobil tarafından gelen hata/uyarı log'ları görünür hale gelecek."),
        ("Temmuz", "Atalian Dashboard Widget",
         "Ay sonu özet kartı: ekstra görev süresi, kayıp sebep dağılımı."),
        ("Ağustos", "Çeklist Şablon Geliştirmesi",
         "Daha esnek şablon, görsel + zorunlu cevap kontrolü."),
        ("Ağustos", "Müşteri Self-Service Rapor",
         "Müşteri kendi rapor şablonunu oluşturabilir, periyodik e-posta gönderimi."),
    ]
    headers_w = [Inches(1.4), Inches(4.5), Inches(7.0)]
    tx, ty = Inches(0.5), Inches(1.4)
    add_box(s, tx, ty, sum(headers_w, Emu(0)), Inches(0.45), fill=ATALIAN_YESIL_K)
    for i, h in enumerate(["DÖNEM", "BAŞLIK", "AÇIKLAMA"]):
        x = tx + sum(headers_w[:i], Emu(0))
        add_text(s, x + Inches(0.1), ty + Inches(0.08), headers_w[i], Inches(0.3),
                 h, font_size=11, bold=True, color=BEYAZ)
    for ri, (donem, baslik, aciklama) in enumerate(items):
        y = ty + Inches(0.45 + 0.7 * ri)
        bg = ACIK_GRI if ri % 2 else BEYAZ
        for ci, val in enumerate([donem, baslik, aciklama]):
            x = tx + sum(headers_w[:ci], Emu(0))
            add_box(s, x, y, headers_w[ci], Inches(0.7), fill=bg,
                    line=PLACEHOLDER_KEN, line_w=Pt(0.5))
            renk = ATALIAN_YESIL_K if ci == 0 else KOYU_GRI
            bold = ci <= 1
            fs   = 11 if ci != 2 else 10
            add_text(s, x + Inches(0.1), y + Inches(0.15), headers_w[ci] - Inches(0.2),
                     Inches(0.5), val, font_size=fs, bold=bold, color=renk)
    return s

def s_kapanis():
    s = add_slide(arka_renk=KOYU_GRI)
    add_box(s, 0, 0, Inches(0.4), prs.slide_height, fill=ATALIAN_YESIL)
    add_text(s, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "TEŞEKKÜRLER", font_size=60, bold=True, color=BEYAZ)
    add_text(s, Inches(1), Inches(3.7), Inches(11), Inches(0.5),
             "Soru & Cevap", font_size=24, color=ATALIAN_ACIK)
    add_box(s, Inches(1), Inches(4.7), Inches(11.3), Inches(0.04),
            fill=ATALIAN_YESIL)
    add_text(s, Inches(1), Inches(5.0), Inches(11), Inches(0.4),
             "ATALIAN — Hizmet Sağlayıcı", font_size=14, bold=True,
             color=OYAK_SARI)
    add_text(s, Inches(1), Inches(5.4), Inches(11), Inches(0.4),
             "OYAK RENAULT — Müşteri",
             font_size=14, color=BEYAZ)
    add_text(s, Inches(1), Inches(6.3), Inches(11), Inches(0.4),
             "İletişim: [Atalian müşteri temsilcisi bilgileri]",
             font_size=11, color=ORTA_GRI)
    add_text(s, Inches(1), Inches(6.7), Inches(11), Inches(0.4),
             "Doküman tarihi: 2026-06",
             font_size=10, color=ORTA_GRI)
    return s

# ── ÜRETİM ───────────────────────────────────────────────────────────────
sayfa_fonlari = [
    s_kapak,
    s_icindekiler,
    s_sistem_genel,
    s_web_yetki,
    s_web_dashboard,
    s_web_canli,
    s_web_kurallar,
    s_mobil_akis,
    s_qr_sure,
    s_ekstra_gorev,
    s_rapor_genel,
    s_rapor_personel,
    s_hakedis,
    s_guvenlik,
    s_musteri_faydalari,
    s_yol_haritasi,
    s_kapanis,
]
total = len(sayfa_fonlari)
slaytlar = []
for i, fn in enumerate(sayfa_fonlari, start=1):
    s = fn()
    slaytlar.append((s, i))

# Footer'ları kapak ve kapanış hariç ekle
for s, i in slaytlar:
    if i in (1, total):  # kapak ve kapanış
        continue
    add_footer(s, i, total)

# Çıktı
out_dir = Path(__file__).resolve().parent.parent.parent / "docs" / "sunum"
out_dir.mkdir(parents=True, exist_ok=True)
out_path = out_dir / "ATALIAN_OYAK_RENAULT_GYS_Sunum.pptx"
prs.save(str(out_path))
print(f"OK: {out_path}")
print(f"Toplam slayt: {total}")
